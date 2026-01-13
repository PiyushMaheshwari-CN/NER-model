# ner_processing.py
import os
import re
import pdfplumber
import docx
import pytesseract
import phonenumbers
import pandas as pd
import spacy
import mysql.connector
from PIL import Image
import pymysql
import snowflake.connector
from pdf2image import convert_from_path

# PATH = r"D:\cloudnexus_tasks\100_resumes"

nlp = spacy.load("en_core_web_sm")

EMAIL_RE = re.compile(r"[a-zA-Z0-9_.+-]+@[a-zA-Z0-9-]+\.[a-zA-Z0-9-.]+")
PHONE_RE = re.compile(r"\+?\d[\d\-\s]{7,}\d")

pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

def extract_text(path):
    ext = os.path.splitext(path)[1].lower()

    try:
        
        if ext == ".pdf":
            text = ""

            
            try:
                with pdfplumber.open(path) as pdf:
                    for p in pdf.pages:
                        page_text = p.extract_text()
                        if page_text:
                            text += "\n" + page_text
            except:
                pass

            # If text is too small, treat it as scanned PDF
            if len(text.strip()) < 20:
                print(f"OCR running for scanned PDF: {path}")
                images = convert_from_path(path)
                for img in images:
                    text += pytesseract.image_to_string(img)

            return text

        # ---------------- DOCX ----------------
        elif ext == ".docx":
            doc = docx.Document(path)
            return "\n".join(p.text for p in doc.paragraphs)

        # ---------------- DOC ----------------
        elif ext == ".doc":
            import textract
            raw = textract.process(path)
            return raw.decode("utf-8", errors="ignore")

        # ---------------- RTF ----------------
        elif ext == ".rtf":
            from striprtf.striprtf import rtf_to_text
            with open(path, "r", errors="ignore") as f:
                return rtf_to_text(f.read())

        # ---------------- ODT ----------------
        elif ext == ".odt":
            import odf.opendocument
            from odf.text import P
            doc = odf.opendocument.load(path)
            paragraphs = doc.getElementsByType(P)
            return "\n".join(p.firstChild.data if p.firstChild else "" for p in paragraphs)

        # ---------------- PPTX ----------------
        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)

        # ---------------- TXT ----------------
        elif ext == ".txt":
            return open(path, "r", encoding="utf-8", errors="ignore").read()

        # ---------------- Images (OCR) ----------------
        elif ext in [".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".webp"]:
            img = Image.open(path)
            return pytesseract.image_to_string(img)

        # ---------------- HTML ----------------
        elif ext == ".html":
            from bs4 import BeautifulSoup
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                soup = BeautifulSoup(f.read(), "html.parser")
                return soup.get_text(separator="\n")

        # ---------------- CSV ----------------
        elif ext == ".csv":
            import pandas as pd
            df = pd.read_csv(path)
            return df.to_string()

        # ---------------- JSON ----------------
        elif ext == ".json":
            import json
            data = json.load(open(path, "r", encoding="utf-8"))
            return json.dumps(data, indent=4)

    except Exception as e:
        print(f"Extraction Error ({path}): {e}")
        return ""

    return ""

def extract_name(text):
    lines = text.splitlines()
    lines = [l.strip() for l in lines if l.strip()][:120]

    BAD_WORDS = {
        "experience", "summary", "profile", "developer", "engineer",
        "skills", "project", "resume", "education", "objective",
        "data", "analyst", "analysis", "curriculum", "intern", "manager",
        "company", "role", "contact", "email", "phone", "mobile",
        "address", "linkedin", "github"
    }

    # ------------------------------------
    # 1️⃣ DIRECT PATTERN: "Name: XYZ"
    # ------------------------------------
    for line in lines:
        if line.lower().startswith(("name:", "name -", "name ")):
            name = line.split(":", 1)[-1].replace("-", "").strip()
            if 1 < len(name.split()) <= 6:
                return name

    # ------------------------------------
    # 2️⃣ TOP BIG FONT DETECTION (First 5 Lines)
    # ------------------------------------
    for line in lines[:5]:
        if 2 <= len(line.split()) <= 4 and not any(c.isdigit() for c in line):
            if all(w.lower() not in BAD_WORDS for w in line.split()):
                return line.strip()

    # ------------------------------------
    # 3️⃣ PUNCTUATION CLEAN + PERSON NER
    # ------------------------------------
    header_text = " ".join(lines[:80])
    header_text = re.sub(r'[^a-zA-Z\s]', ' ', header_text)

    doc = nlp(header_text)
    for ent in doc.ents:
        if ent.label_ == "PERSON":
            name = ent.text.strip()
            if (1 < len(name.split()) <= 6 
                and all(w.lower() not in BAD_WORDS for w in name.split())):
                return name

    # ------------------------------------
    # 4️⃣ FALLBACK: FIRST CLEAN HUMAN-LIKE LINE
    # ------------------------------------
    for line in lines[:40]:
        if any(x in line.lower() for x in BAD_WORDS):
            continue
        if re.search(r"[0-9@#/$%&*<>]", line):
            continue
        parts = line.split()
        if 1 < len(parts) <= 4:
            return line.strip()

    # ------------------------------------
    # 5️⃣ LAST FALLBACK
    # ------------------------------------
    return "NA"



def extract_email(text):
    found = EMAIL_RE.findall(text)
    return found[0] if found else "NA"


def extract_phone(text):
    clean = re.sub(r"[^0-9+]", " ", text)
    tokens = clean.split()
    candidates = []
    for token in tokens:
        num = re.sub(r"[^\d+]", "", token)
        if len(num) >= 8:
            candidates.append(num)

    if not candidates:
        return "NA"

    for num in candidates:
        try:
            parsed = phonenumbers.parse(num, "IN")
            if phonenumbers.is_valid_number(parsed):
                return phonenumbers.format_number(parsed, phonenumbers.PhoneNumberFormat.E164)
        except:
            pass

    longest = max(candidates, key=len)
    if longest.startswith("91") and len(longest) == 12:
        longest = "+" + longest
    elif len(longest) == 10:
        longest = "+91" + longest

    return longest


# def extract_skills(text):
#     """
#     Extract ONLY technical skills from the resume's SKILL section.
#     Uses:
#       - Section detection (Skills, Technical Skills, Key Skills)
#       - Pattern-based technical filtering
#     """

#     if not text:
#         return ["NA"]

#     text_lower = text.lower()

#     section_patterns = [
#         r"skills\s*[:\-]?", 
#         r"technical skills\s*[:\-]?",
#         r"key skills\s*[:\-]?",
#         r"skills & tools\s*[:\-]?",
#         r"expertise\s*[:\-]?"
#     ]

#     skill_section = ""

#     for pattern in section_patterns:
#         match = re.search(pattern, text_lower)
#         if match:
#             start = match.end()
#             skill_section = text_lower[start:start + 400]
#             break

#     if not skill_section:
#         return ["NA"]   

#     skill_section = re.split(
#         r"(education|project|projects|experience|work|summary)\s*[:\-]?",
#         skill_section
#     )[0]

#     cleaned = re.sub(r"[^a-z0-9+.#/\s-]", " ", skill_section)
#     tokens = re.split(r"[\s,;:/\|\-\n]+", cleaned)

#     skills = set()

#     prog_lang = r"^(python|java|javascript|typescript|c|c\+\+|c#|go|php|ruby|swift|kotlin)$"

#     frameworks = r"^(react|next\.js|node\.js|express|django|flask|spring|angular|vue|pandas|numpy|matplotlib|pytorch|tensorflow|sklearn|selenium)$"

#     tools = r"^(aws|azure|gcp|docker|kubernetes|jenkins|git|github|gitlab|terraform)$"

#     db = r"^(mysql|postgres|postgresql|sqlite|oracle|mongodb|redis|firebase|mssql)$"

#     def dynamic_rule(t):
#         return (
#             any(x in t for x in ["+", "#", ".", "js", "sql", "api"]) or
#             any(c.isdigit() for c in t) or
#             (t.isupper() and 1 < len(t) <= 6)
#         )

#     for tok in tokens:
#         tok = tok.strip()
#         if len(tok) < 2:
#             continue

#         if re.match(prog_lang, tok):
#             skills.add(tok)
#             continue

#         if re.match(frameworks, tok):
#             skills.add(tok)
#             continue

#         if re.match(tools, tok):
#             skills.add(tok)
#             continue

#         if re.match(db, tok):
#             skills.add(tok)
#             continue

#         if dynamic_rule(tok):
#             skills.add(tok)

#     final = []
#     for s in skills:
#         if any(x in s for x in ["+", "#", "."]):
#             final.append(s)
#         else:
#             final.append(s.title())

#     return final if final else ["NA"]

def extract_skills(text):
    if not text:
        return ["NA"]

    text_lower = text.lower()

    # --------------------------
    # 1. Try finding SKILLS section
    # --------------------------
    section_patterns = [
        r"skills\s*[:\-]?",
        r"technical skills\s*[:\-]?",
        r"key skills\s*[:\-]?",
        r"skills & tools\s*[:\-]?",
        r"expertise\s*[:\-]?"
    ]

    skill_section = ""

    for pattern in section_patterns:
        match = re.search(pattern, text_lower)
        if match:
            start = match.end()
            skill_section = text_lower[start:start + 400]
            break

    # If not found, fallback to whole resume
    if not skill_section:
        skill_section = text_lower

    # --------------------------
    # 2. Remove unwanted text
    # --------------------------
    skill_section = re.split(
        r"(education|project|projects|experience|work|summary)\s*[:\-]?",
        skill_section
    )[0]

    cleaned = re.sub(r"[^a-z0-9+.#/\s-]", " ", skill_section)
    tokens = re.split(r"[\s,;:/\|\-\n]+", cleaned)

    skills = set()

    # --------------------------
    # 3. Dictionaries
    # --------------------------
    prog_lang = r"^(python|java|javascript|typescript|c|c\+\+|c#|go|php|ruby|swift|kotlin)$"
    frameworks = r"^(react|next\.js|node\.js|express|django|flask|spring|angular|vue|pandas|numpy|matplotlib|pytorch|tensorflow|sklearn|selenium)$"
    tools = r"^(aws|azure|gcp|docker|kubernetes|jenkins|git|github|gitlab|terraform)$"
    db = r"^(mysql|postgres|postgresql|sqlite|oracle|mongodb|redis|firebase|mssql)$"

    # --------------------------
    # 4. Dynamic intelligent rules
    # --------------------------
    def dynamic_rule(t):
        return (
            any(x in t for x in ["+", "#", ".", "js", "sql", "api"]) or
            any(c.isdigit() for c in t) or
            (t.isupper() and 1 < len(t) <= 6)
        )

    # --------------------------
    # 5. Main extraction
    # --------------------------
    for tok in tokens:
        tok = tok.strip()
        if len(tok) < 2:
            continue

        if re.match(prog_lang, tok): skills.add(tok); continue
        if re.match(frameworks, tok): skills.add(tok); continue
        if re.match(tools, tok): skills.add(tok); continue
        if re.match(db, tok): skills.add(tok); continue
        if dynamic_rule(tok): skills.add(tok)

    final = [s if any(x in s for x in ["+", "#", "."]) else s.title() for s in skills]

    return final if final else ["NA"]



def get_section(text, keys):
    lines = text.splitlines()
    capture = False
    result = []

    for ln in lines:
        low = ln.lower()

        if any(k in low for k in keys):
            capture = True
            continue

        if capture and any(stop in low for stop in
            ["education","experience","summary","contact","project","skills"]):
            break

        if capture:
            result.append(ln.strip())

    return result


def extract_projects(text):
    section = get_section(text, ["project"])
    if section:
        cleaned = [p.strip() for p in section if len(p.strip()) > 5]
        return cleaned if cleaned else ["NA"]

    found = [l.strip() for l in text.splitlines() if "project" in l.lower()]
    return found if found else ["NA"]


def extract_resume_data(file_path):
    text = extract_text(file_path)

    return {
        "name": extract_name(text),
        "email": extract_email(text),
        "phone": extract_phone(text),
        "skills": extract_skills(text),
        "projects": extract_projects(text)
    }


def get_all_files(path):
    if os.path.isfile(path):
        return [path]
    else:
        result = []
        for f in os.listdir(path):
            full = os.path.join(path, f)
            if os.path.isfile(full):
                result.append(full)
        return result


def process_resumes(file_paths):
    """
    Called by Flask UI.
    file_paths: list of saved uploaded file paths (in uploads/)
    Runs full pipeline: extract -> CSVs -> MySQL insert -> export temp.csv -> S3 upload -> Snowflake load
    """
    names, emails, phones, skills, projects = [], [], [], [], []

    # build all_files from provided file_paths
    all_files = [f for f in file_paths if os.path.isfile(f)]

    for path in all_files:
        file = os.path.basename(path)
        ext = os.path.splitext(file)[1].lower()

        if ext not in [".pdf", ".docx", ".doc", ".txt", ".jpg", ".jpeg", ".png", ".rtf",".json",".csv",".html",".bmp", ".tiff", ".webp",".txt",".pptx",".odt"]:
            print(f"Skipping unsupported file: {file}")
            continue

        print(f"Processing: {file}")

        text = extract_text(path)

        name_val = extract_name(text)
        email_val = extract_email(text)
        phone_val = extract_phone(text)
        file_skills = extract_skills(text)
        file_projects = extract_projects(text)

        names.append({"file": file, "name": name_val})
        emails.append({"file": file, "email": email_val})
        phones.append({"file": file, "phone": phone_val})

        for s in file_skills:
            skills.append({"file": file, "skill": s})

        for p in file_projects:
            projects.append({"file": file, "project": p})

    # Write CSV files (same names as before)
    pd.DataFrame(names).to_csv("name.csv", index=False)
    pd.DataFrame(emails).to_csv("email.csv", index=False)
    pd.DataFrame(phones).to_csv("phone.csv", index=False)
    pd.DataFrame(skills).to_csv("skills.csv", index=False)
    pd.DataFrame(projects).to_csv("projects.csv", index=False)

    print("\n✔✔ All CSV files generated successfully!")

    # -------------------------
    # MySQL insertion
    # -------------------------
    names_df = pd.read_csv("name.csv")
    emails_df = pd.read_csv("email.csv")
    phones_df = pd.read_csv("phone.csv")
    skills_df = pd.read_csv("skills.csv")
    projects_df = pd.read_csv("projects.csv")

    df = names_df.merge(emails_df, on="file", how="left") \
                 .merge(phones_df, on="file", how="left")

    # MySQL connection (same credentials as you had)
    db = mysql.connector.connect(
        host="localhost",
        user="root",
        password="mysql",
        database="ner_1"
    )
    cursor = db.cursor()

    insert_query = """
    INSERT INTO resume_data (name, email, phone, skill, project)
    VALUES (%s, %s, %s, %s, %s)
    """

    for file in df["file"].unique():
        name = df[df["file"] == file]["name"].iloc[0]
        email = df[df["file"] == file]["email"].iloc[0]
        phone = df[df["file"] == file]["phone"].iloc[0]

        phone2 = str(phone)
        file_skills = skills_df[skills_df["file"] == file]["skill"].tolist()
        file_projects = projects_df[projects_df["file"] == file]["project"].tolist()

        if not file_skills:
            file_skills = [None]
        if not file_projects:
            file_projects = [None]

        # for s in file_skills:
        #     for p in file_projects:
        cursor.execute(insert_query, (name, email, phone2, s, p))

    db.commit()
    cursor.close()
    db.close()

    print("✔ Data inserted into MySQL successfully!")

    # -------------------------
    # Export MySQL table to temp.csv
    # -------------------------
    mysql_host = "localhost"
    mysql_user = "root"
    mysql_password = "mysql"
    mysql_database = "ner_1"
    mysql_table = "resume_data"

    conn = pymysql.connect(
        host=mysql_host,
        user=mysql_user,
        password=mysql_password,
        database=mysql_database
    )

    df_sql = pd.read_sql(f"SELECT * FROM {mysql_table}", conn)
    print("MySQL Table Loaded Successfully!")

    df_sql.to_csv("temp.csv", index=False)
    print("Converted to CSV!")

    conn.close()

    # -------------------------
    # Upload to S3 (same variables as in your original code)
    # -------------------------
    # original script had AWS variables — keep them as-is here
    AWS_ACCESS_KEY ="AKIAWV2KVQMYQ5IIUF66"
    AWS_SECRET_KEY ="dMivI1ryjW9NhGmWyi1HOjadl1ly0v4oB6QCT94w"
    BUCKET_NAME = "piyush.20041"
    FILE_NAME = "resume_data.csv"
    REGION = "ap-southeast-2"

    # if keys are provided, attempt upload
    if AWS_ACCESS_KEY and AWS_SECRET_KEY:
        s3 = boto3.client(
            "s3",
            aws_access_key_id = AWS_ACCESS_KEY,
            aws_secret_access_key = AWS_SECRET_KEY,
            region_name = REGION
        )

        try:
            s3.upload_file(
                Filename="temp.csv",
                Bucket=BUCKET_NAME,
                Key=FILE_NAME
            )
            print("✔ File Uploaded to S3 Successfully!")
        except Exception as e:
            print("Error uploading:", e)
    else:
        print("Skipping S3 upload (no AWS keys provided in code).")

    # -------------------------
    # Snowflake load (same as original)
    # -------------------------
    SNOWFLAKE_USER = "Piyush2004"
    SNOWFLAKE_PASSWORD = "Piyush1maheshwari"
    SNOWFLAKE_ACCOUNT = "KFNNGJM-ZP92760"
    SNOWFLAKE_WAREHOUSE = "COMPUTE_WH"
    SNOWFLAKE_DATABASE = "CONNECT_BI"
    SNOWFLAKE_SCHEMA = "PUBLIC"

    S3_BUCKET = "piyush.20041"
    S3_FILE_KEY = "https://s3.ap-southeast-2.amazonaws.com/piyush.20041/resume_data.csv"
    AWS_ACCESS_KEY = "AKIAWV2KVQMYQ5IIUF66"
    AWS_SECRET_KEY = "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"

    try:
        conn_sf = snowflake.connector.connect(
            user=SNOWFLAKE_USER,
            password=SNOWFLAKE_PASSWORD,
            account=SNOWFLAKE_ACCOUNT,
            warehouse=SNOWFLAKE_WAREHOUSE,
            database=SNOWFLAKE_DATABASE,
            schema=SNOWFLAKE_SCHEMA
        )

        cursor = conn_sf.cursor()

        cursor.execute("""
        CREATE TABLE IF NOT EXISTS resume_data (
            ID INTEGER,
            name STRING,
            email STRING,
            phone STRING,
            skill STRING,
            project STRING
        );
        """)
        print("Table checked/created.")

        cursor.execute(f"""
        COPY INTO resume_data
        FROM @S3_STAGE/{S3_FILE_KEY}
        FILE_FORMAT = (TYPE = 'CSV' FIELD_OPTIONALLY_ENCLOSED_BY='\"' SKIP_HEADER = 1)
        ON_ERROR = 'CONTINUE';
        """)
        print("CSV loaded into Snowflake from existing stage!")

        cursor.close()
        conn_sf.close()
    except Exception as e:
        print("Snowflake step failed (check credentials / stage):", e)

    return True
